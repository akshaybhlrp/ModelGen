import os
import sys
import threading
import concurrent.futures
from pathlib import Path
from kernel import init_db, store
from harvester import extract_ast
from mutation_tester import evaluate_mutation_score
from decontaminate import DecontaminationGate
import subprocess
import shutil

RTK_BIN = shutil.which("rtk") or "/home/akshay-bhalerao/.local/bin/rtk"

def rtk_smart_summarize(file_path: Path) -> str:
    """Uses Rust Token Killer (RTK) binary for high-speed AST / heuristic token-compressed file summaries."""
    if os.path.exists(RTK_BIN):
        try:
            r = subprocess.run([RTK_BIN, "smart", str(file_path)], capture_output=True, text=True, timeout=2)
            if r.returncode == 0 and r.stdout.strip():
                return r.stdout.strip()
        except Exception:
            pass
    return ""

def rtk_tree_structure(dir_path: Path) -> str:
    """Uses RTK to generate a compact, token-optimized directory tree structure."""
    if os.path.exists(RTK_BIN):
        try:
            r = subprocess.run([RTK_BIN, "tree", str(dir_path), "--max-depth", "3"], capture_output=True, text=True, timeout=3)
            if r.returncode == 0 and r.stdout.strip():
                return r.stdout.strip()
        except Exception:
            pass
    return ""

# Universal Multi-Format Matrix
SUPPORTED_CODE_EXTS = {".py", ".ipynb", ".sh", ".bash", ".js", ".ts", ".c", ".cpp", ".h", ".rs", ".go", ".java", ".sql"}
SUPPORTED_OFFICE_DOCS = {".docx", ".doc", ".odt", ".rtf", ".pdf", ".epub"}
SUPPORTED_SPREADSHEETS = {".xlsx", ".xls", ".csv", ".tsv", ".ods"}
SUPPORTED_PRESENTATIONS = {".pptx", ".ppt", ".odp"}
SUPPORTED_CAD_3D_EXTS = {".dxf", ".dwg", ".step", ".stp", ".iges", ".igs", ".stl", ".obj", ".gltf", ".glb", ".fbx", ".ply", ".3ds", ".dae"}
SUPPORTED_DATA_CONFIG_EXTS = {".json", ".yaml", ".yml", ".xml", ".toml", ".ini", ".cfg", ".md", ".txt", ".html", ".htm"}
SUPPORTED_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".svg", ".tiff", ".gif", ".ico"}
SUPPORTED_VIDEO_EXTS = {".mp4", ".mov", ".mkv", ".avi", ".webm", ".flv", ".wmv"}
SUPPORTED_AUDIO_EXTS = {".mp3", ".wav", ".m4a", ".ogg", ".flac", ".aac"}
SUPPORTED_ARCHIVE_EXTS = {".zip", ".tar", ".gz", ".tgz", ".bz2", ".tar.gz", ".tar.bz2", ".7z", ".rar", ".xz"}

ALL_MULTIMODAL_EXTS = (
    SUPPORTED_OFFICE_DOCS | SUPPORTED_SPREADSHEETS | SUPPORTED_PRESENTATIONS |
    SUPPORTED_CAD_3D_EXTS | SUPPORTED_DATA_CONFIG_EXTS | SUPPORTED_IMAGE_EXTS |
    SUPPORTED_VIDEO_EXTS | SUPPORTED_AUDIO_EXTS
)

def extract_office_text(file_path: Path) -> str:
    """Extracts text content and tables from Microsoft Office / OpenDocument formats."""
    ext = file_path.suffix.lower()
    text = []
    
    # 1. Word Documents (.docx)
    if ext == ".docx":
        try:
            import docx
            doc = docx.Document(str(file_path))
            for p in doc.paragraphs:
                if p.text.strip():
                    text.append(p.text.strip())
        except Exception:
            pass

    # 2. Excel Spreadsheets (.xlsx, .csv, .tsv)
    elif ext in {".xlsx", ".xls", ".ods"}:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
            for sheet in wb.sheetnames[:5]:
                ws = wb[sheet]
                text.append(f"--- Sheet: {sheet} ---")
                for row in ws.iter_rows(max_row=20, values_only=True):
                    row_vals = [str(v) for v in row if v is not None]
                    if row_vals:
                        text.append(" | ".join(row_vals))
        except Exception:
            pass

    # 3. PowerPoint Presentations (.pptx)
    elif ext == ".pptx":
        try:
            import pptx
            prs = pptx.Presentation(str(file_path))
            for idx, slide in enumerate(prs.slides[:30]):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_texts.append(shape.text.strip())
                if slide_texts:
                    text.append(f"[Slide {idx+1}]: " + " - ".join(slide_texts))
        except Exception:
            pass

    return "\n".join(text) if text else f"Office Document: {file_path.name}"

def extract_cad_metadata(cad_path: Path) -> str:
    """Parses 2D/3D CAD geometry, layers, entities, and mesh statistics from CAD formats."""
    ext = cad_path.suffix.lower()
    desc = [f"CAD / 3D Asset: {cad_path.name} ({ext[1:].upper()})"]
    size_kb = round(cad_path.stat().st_size / 1024, 2)
    desc.append(f"File Size: {size_kb} KB")
    
    # 1. DXF CAD Vector Parsing
    if ext == ".dxf":
        try:
            import ezdxf
            doc = ezdxf.readfile(str(cad_path))
            msp = doc.modelspace()
            layers = [layer.dxf.name for layer in doc.layers]
            desc.append(f"CAD Layers ({len(layers)}): {', '.join(layers[:8])}")
            desc.append(f"ModelSpace Entities Count: {len(msp)}")
        except Exception:
            # Fallback simple header inspection
            try:
                txt = cad_path.read_text(errors="ignore")[:500]
                if "$ACADVER" in txt:
                    desc.append("AutoCAD DXF Header Validated")
            except Exception:
                pass

    # 2. STL / OBJ 3D Meshes
    elif ext in {".stl", ".obj", ".ply", ".gltf", ".glb"}:
        try:
            # Scan vertex/face lines in text OBJ/STL
            content = cad_path.read_text(errors="ignore")[:2000]
            v_count = content.count("\nv ")
            f_count = content.count("\nf ")
            if v_count or f_count:
                desc.append(f"Estimated Mesh Sample: {v_count} vertices, {f_count} faces")
        except Exception:
            pass

    # 3. STEP / IGES Neutral CAD Standards
    elif ext in {".step", ".stp", ".iges", ".igs"}:
        try:
            header_lines = []
            with open(cad_path, "r", errors="ignore") as f:
                for _ in range(25):
                    line = f.readline()
                    if not line: break
                    if "HEADER;" in line or "FILE_NAME" in line or "FILE_SCHEMA" in line or "Start Section" in line:
                        header_lines.append(line.strip())
            if header_lines:
                desc.append("CAD Standards Header:\n" + "\n".join(header_lines[:5]))
        except Exception:
            pass

    return "\n".join(desc)

def extract_pdf_text(pdf_path: Path) -> str:
    """Extracts text, code blocks, and doc contents from a PDF file."""
    text_content = []
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(pdf_path))
        for page in reader.pages[:50]:
            t = page.extract_text()
            if t:
                text_content.append(t)
    except Exception:
        pass
    
    if not text_content:
        try:
            import subprocess
            r = subprocess.run(["pdftotext", str(pdf_path), "-"], capture_output=True, text=True, timeout=5)
            if r.returncode == 0 and r.stdout:
                text_content.append(r.stdout)
        except Exception:
            pass

    return "\n".join(text_content) if text_content else f"PDF Document: {pdf_path.name}"

def process_archive_file(archive_path: Path, conn) -> int:
    """Unpacks compressed archives in an isolated temp folder and ingests contained code & media."""
    learned = 0
    with tempfile.TemporaryDirectory() as tmp_extract_dir:
        tmp_target = Path(tmp_extract_dir)
        try:
            if zipfile.is_zipfile(archive_path):
                with zipfile.ZipFile(archive_path, 'r') as zf:
                    zf.extractall(tmp_target)
            elif tarfile.is_tarfile(archive_path):
                with tarfile.open(archive_path, 'r:*') as tf:
                    tf.extractall(tmp_target)
            
            sub_res = ingest_local_directory(str(tmp_target), conn=conn, retrain_neural_weights=False, max_workers=4)
            learned = sub_res.get("learned_count", 0)
        except Exception:
            pass
    return learned

def process_media_file(file_path: Path, conn) -> list:
    """Extracts algorithmic patterns, OCR text, CAD vectors, or specs from any document or multimodal file."""
    items = []
    ext = file_path.suffix.lower()
    
    # 1. PDF Documents
    if ext == ".pdf":
        try:
            pdf_text = extract_pdf_text(file_path)
            from harvester import extract_ast
            extracted_code = extract_ast(pdf_text)
            if extracted_code:
                for fn_name, fn_code, test_code in extracted_code:
                    items.append((fn_name, fn_code, test_code))
            else:
                items.append((f"pdf_{file_path.stem}", f"# PDF Document {file_path.name}\n# Content Summary:\n# {pdf_text[:400]}\npass", "def test():\n    pass\n"))
        except Exception:
            pass

    # 2. Office Documents (Word, Excel, PowerPoint)
    elif ext in (SUPPORTED_OFFICE_DOCS | SUPPORTED_SPREADSHEETS | SUPPORTED_PRESENTATIONS):
        try:
            doc_text = extract_office_text(file_path)
            from harvester import extract_ast
            extracted_code = extract_ast(doc_text)
            if extracted_code:
                for fn_name, fn_code, test_code in extracted_code:
                    items.append((fn_name, fn_code, test_code))
            else:
                items.append((f"office_{file_path.stem}", f"# Office Asset: {file_path.name}\n# Summary:\n# {doc_text[:400]}\npass", "def test():\n    pass\n"))
        except Exception:
            pass

    # 3. CAD & 3D Vector Formats (.dxf, .step, .stl, .obj, .dwg, .gltf)
    elif ext in SUPPORTED_CAD_3D_EXTS:
        try:
            cad_desc = extract_cad_metadata(file_path)
            items.append((f"cad_{file_path.stem}", f"# CAD Geometry Asset: {file_path.name}\n# Specification:\n# {cad_desc}\npass", "def test():\n    pass\n"))
        except Exception:
            pass

    # 4. Images (OCR or image captioning extraction)
    elif ext in SUPPORTED_IMAGE_EXTS:
        try:
            from PIL import Image
            img = Image.open(file_path)
            w, h = img.size
            desc = f"Image asset: {file_path.name} (dimensions {w}x{h})"
            try:
                import pytesseract
                ocr_text = pytesseract.image_to_string(img).strip()
                if ocr_text:
                    desc += f"\nOCR Extracted text:\n{ocr_text}"
            except Exception:
                pass
            items.append((f"img_{file_path.stem}", f"# {desc}\npass", "def test():\n    pass\n"))
        except Exception:
            pass

    # 5. Videos / Audio (metadata and transcript hooks)
    elif ext in SUPPORTED_VIDEO_EXTS or ext in SUPPORTED_AUDIO_EXTS:
        try:
            size_mb = round(file_path.stat().st_size / (1024 * 1024), 2)
            desc = f"Media file {file_path.name} ({ext[1:].upper()}, {size_mb} MB)"
            items.append((f"media_{file_path.stem}", f"# {desc}\npass", "def test():\n    pass\n"))
        except Exception:
            pass

    # 6. Data & Config specs (.json, .yaml, .csv, .xml, .html)
    elif ext in SUPPORTED_DATA_CONFIG_EXTS:
        try:
            txt = file_path.read_text(errors="ignore")[:600]
            items.append((f"doc_{file_path.stem}", f"# Config/Data Asset: {file_path.name}\n# Header:\n# {txt}\npass", "def test():\n    pass\n"))
        except Exception:
            pass

    # 7. UNIVERSAL CATCH-ALL FOR ANY UNLISTED / CUSTOM / BINARY / SCIENTIFIC FILE FORMAT
    # (e.g. .h5, .parquet, .feather, .fits, .nii, .dcm, .wasm, .dat, .bin, .log, .raw, .mat, etc.)
    else:
        try:
            size_kb = round(file_path.stat().st_size / 1024, 2)
            # Try reading as text first (for unknown text/logs/source files)
            try:
                sample_txt = file_path.read_text(encoding="utf-8", errors="strict")[:500].strip()
                if sample_txt and sample_txt.isprintable():
                    from harvester import extract_ast
                    extracted_code = extract_ast(sample_txt)
                    if extracted_code:
                        for fn_name, fn_code, test_code in extracted_code:
                            items.append((fn_name, fn_code, test_code))
                    else:
                        items.append((f"data_{file_path.stem}", f"# Custom Text Asset: {file_path.name}\n# Sample Content:\n# {sample_txt[:300]}\npass", "def test():\n    pass\n"))
                else:
                    items.append((f"blob_{file_path.stem}", f"# Binary Asset: {file_path.name} (Format: {ext or 'raw'}, Size: {size_kb} KB)\npass", "def test():\n    pass\n"))
            except Exception:
                # Pure binary file fallback (extract header bytes & magic metadata)
                with open(file_path, "rb") as bf:
                    raw_head = bf.read(64).hex()
                items.append((f"bin_{file_path.stem}", f"# Binary Asset: {file_path.name} ({size_kb} KB)\n# Header Hex: {raw_head[:32]}\npass", "def test():\n    pass\n"))
        except Exception:
            pass

    return items

# Global Real-Time Scan Progress State
SCAN_PROGRESS = {
    "active": False,
    "directory": "",
    "total_files": 0,
    "processed_files": 0,
    "learned_count": 0,
    "current_file": "",
    "percent": 0,
    "status": "idle"
}
PROGRESS_LOCK = threading.Lock()

def get_scan_progress():
    with PROGRESS_LOCK:
        return dict(SCAN_PROGRESS)

def set_scan_progress(**kwargs):
    with PROGRESS_LOCK:
        SCAN_PROGRESS.update(kwargs)
        if SCAN_PROGRESS["total_files"] > 0:
            SCAN_PROGRESS["percent"] = min(100, int((SCAN_PROGRESS["processed_files"] / SCAN_PROGRESS["total_files"]) * 100))

def ingest_local_directory(dir_path: str, conn=None, retrain_neural_weights: bool = True, max_workers: int = 8, progress_callback=None) -> dict:
    """Asynchronous, multi-worker recursive directory scanner for code and multimodal media."""
    target = Path(dir_path).resolve()
    if not target.exists() or not target.is_dir():
        return {"status": "error", "message": f"Directory not found: {dir_path}", "learned_count": 0}

    if conn is None:
        conn = init_db()

    decontam = DecontaminationGate()
    
    # Discover all code, media, and docs recursively
    all_files = []
    for root, dirs, files in os.walk(target):
        # Ignore hidden / build / virtualenv directories
        dirs[:] = [d for d in dirs if not d.startswith(".") and d not in {"venv", ".venv", "__pycache__", "node_modules", "build", "dist"}]
        for f in files:
            # Ignore hidden files, SQLite databases, and WAL logs
            if not f.startswith(".") and not f.endswith(".db") and not f.endswith("-wal") and not f.endswith("-shm"):
                all_files.append(Path(root) / f)

    code_files = [f for f in all_files if f.suffix.lower() in SUPPORTED_CODE_EXTS]
    archive_files = [f for f in all_files if any(f.name.lower().endswith(ext) for ext in SUPPORTED_ARCHIVE_EXTS)]
    # All remaining files (office, CAD, media, database, scientific, binary, custom) route to universal ingestion
    media_files = [f for f in all_files if f not in code_files and f not in archive_files]

    total_found = 0
    total_stored = 0
    modules_stored = []
    processed_count = 0
    lock = threading.Lock()

    set_scan_progress(
        active=True,
        directory=str(target.name),
        total_files=len(all_files),
        processed_files=0,
        learned_count=0,
        current_file="Initializing scan...",
        status="scanning"
    )

    print(f"\n[+] Async Scanning Directory: {target} ({len(code_files)} Code, {len(media_files)} All Multi-Format Files, {len(archive_files)} Archives across {max_workers} threads)")

    def update_file_progress(file_path: Path):
        nonlocal processed_count
        with lock:
            processed_count += 1
            set_scan_progress(
                processed_files=processed_count,
                current_file=file_path.name,
                learned_count=total_stored
            )
            if progress_callback:
                progress_callback(processed_count, len(all_files), file_path.name, total_stored)

    def process_single_code_file(py_file: Path):
        nonlocal total_found, total_stored
        active_conn = conn if conn is not None else init_db()
        try:
            content = py_file.read_text(encoding="utf-8", errors="ignore")
            if len(content) < 30:
                return
            extracted = extract_ast(content)
            for fn_name, fn_code, test_code in extracted:
                with lock:
                    total_found += 1

                # Decontamination check
                is_contam, _ = decontam.is_contaminated(fn_code, test_code, source_url=f"local:{py_file}")
                if is_contam:
                    continue

                # Mutation quality check
                mut_score, killed, total = evaluate_mutation_score(fn_code, test_code, max_mutants=5)
                if mut_score < 0.40:
                    continue

                with lock:
                    mid = store(active_conn, fn_name, fn_code, test_code, "Local", f"local:{py_file}")
                    if mid:
                        total_stored += 1
                        modules_stored.append(fn_name)
                        print(f"  [VERIFIED & LEARNED] #{mid} '{fn_name}' from {py_file.name} (Mutation Kill-Rate: {mut_score:.1%})")
        except Exception:
            pass
        finally:
            update_file_progress(py_file)

    def process_single_media(m_file: Path):
        nonlocal total_found, total_stored
        active_conn = conn if conn is not None else init_db()
        try:
            m_items = process_media_file(m_file, active_conn)
            for m_name, m_code, m_tests in m_items:
                with lock:
                    total_found += 1
                    mid = store(active_conn, m_name, m_code, m_tests, "MediaAsset", f"local_media:{m_file}")
                    if mid:
                        total_stored += 1
                        modules_stored.append(m_name)
                        print(f"  [MEDIA INDEXED] #{mid} '{m_name}' ({m_file.suffix})")
        except Exception:
            pass
        finally:
            update_file_progress(m_file)

    # Process files
    if max_workers <= 1 or len(code_files) + len(media_files) <= 1:
        for f in code_files:
            process_single_code_file(f)
        for f in media_files:
            process_single_media(f)
    else:
        # Multi-threaded concurrent execution pool
        with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
            code_futures = [executor.submit(process_single_code_file, f) for f in code_files]
            media_futures = [executor.submit(process_single_media, f) for f in media_files]
            concurrent.futures.wait(code_futures + media_futures)

    # Process compressed archives
    for arc_file in archive_files:
        try:
            arc_learned = process_archive_file(arc_file, conn)
            total_stored += arc_learned
            if arc_learned > 0:
                print(f"  [ARCHIVE UNPACKED & LEARNED] {arc_file.name} -> {arc_learned} items")
        except Exception:
            pass
        finally:
            update_file_progress(arc_file)

    # On-the-fly neural router weight retraining
    if total_stored > 0 and retrain_neural_weights:
        set_scan_progress(current_file="Adapting neural router weights (InfoNCE)...", status="retraining")
        print(f"\n[+] Adapting Neural Router Weights for {total_stored} newly learned items...")
        train_learned_router(conn, epochs=10)

    set_scan_progress(
        active=False,
        processed_files=len(all_files),
        percent=100,
        current_file="Scan & Ingestion Complete",
        status="complete",
        learned_count=total_stored
    )

    return {
        "status": "success",
        "scanned_files": len(all_files),
        "code_files": len(code_files),
        "media_files": len(media_files),
        "candidates_found": total_found,
        "learned_count": total_stored,
        "modules": modules_stored
    }

def ingest_async_background(dir_path: str, conn=None):
    """Launches directory scan non-blockingly in a background thread."""
    t = threading.Thread(target=ingest_local_directory, args=(dir_path, conn), daemon=True)
    t.start()
    return t

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python3 local_learner.py <path_to_directory>")
        sys.exit(1)

    path = sys.argv[1]
    res = ingest_local_directory(path)
    print(f"\n[+] Async Ingestion Complete. Learned {res['learned_count']} modules from {res['scanned_files']} files ({res['media_files']} media).")
